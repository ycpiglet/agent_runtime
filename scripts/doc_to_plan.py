"""Document -> plan pipeline (TASK-AR-366, TASKSET-AR-DOC-TO-PLAN).

Ingest an idea/pitch/spec document and propose a registerable plan: analyze it into
goals/features/constraints, draft a taskset + task decomposition, and emit a B-mode
PLANNING PROPOSAL (origin_type=doc_intake, status=proposed) shaped for the existing
`work.py new` registration path. Registration is **Owner-gated** — this tool never
writes the registry; it produces an approvable proposal.

Parsing: md/txt/html are stdlib. Binary decks (PDF/PPTX/DOCX) use lazy OPTIONAL local
libraries; if a library is absent the tool reports it gracefully instead of crashing
(the repo + CI are dependency-light). No external services.
"""
from __future__ import annotations

import argparse
import json
import re
from html.parser import HTMLParser
from pathlib import Path

GOAL_HINTS = ("goal", "objective", "vision", "목표", "비전", "목적")
CONSTRAINT_HINTS = ("constraint", "must not", "limit", "제약", "금지", "non-goal")
FEATURE_HINTS = ("feature", "capability", "requirement", "기능", "요구", "해야")
_BULLET = re.compile(r"^\s*(?:[-*+]|\d+[.)])\s+(.*)$")
_HEADING = re.compile(r"^\s*#{1,6}\s+(.*)$")


class _Stripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data):
        if data.strip():
            self.parts.append(data.strip())


def extract_text(path: Path) -> str:
    """Normalize a document to plain text. Stdlib for md/txt/html; lazy for binaries."""
    suffix = path.suffix.lower()
    if suffix in (".md", ".markdown", ".txt", ""):
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix in (".html", ".htm"):
        s = _Stripper()
        s.feed(path.read_text(encoding="utf-8", errors="replace"))
        return "\n".join(s.parts)
    if suffix == ".pdf":
        return _lazy_binary(path, "pdfplumber", _pdf_text)
    if suffix in (".pptx",):
        return _lazy_binary(path, "pptx", _pptx_text)
    if suffix in (".docx",):
        return _lazy_binary(path, "docx", _docx_text)
    raise ValueError(f"unsupported document type: {suffix or '(none)'}")


def _lazy_binary(path: Path, module: str, fn) -> str:
    try:
        __import__(module)
    except ImportError as exc:
        raise RuntimeError(
            f"{path.suffix} needs the optional local library '{module}' "
            f"(install it; no external service): {exc}"
        ) from exc
    return fn(path)


def _pdf_text(path: Path) -> str:  # pragma: no cover - exercised only when pdfplumber present
    import pdfplumber
    with pdfplumber.open(str(path)) as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)


def _pptx_text(path: Path) -> str:  # pragma: no cover
    from pptx import Presentation
    prs = Presentation(str(path))
    return "\n".join(sh.text for slide in prs.slides for sh in slide.shapes if hasattr(sh, "text"))


def _docx_text(path: Path) -> str:  # pragma: no cover
    import docx
    return "\n".join(p.text for p in docx.Document(str(path)).paragraphs)


def analyze(text: str, *, title: str | None = None) -> dict:
    goals, features, constraints, milestones = [], [], [], []
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        h = _HEADING.match(raw)
        if h:
            head = h.group(1).strip()
            if not title:
                title = head
            else:
                milestones.append(head)
            continue
        b = _BULLET.match(raw)
        content = b.group(1).strip() if b else line
        low = content.lower()
        if any(k in low for k in CONSTRAINT_HINTS):
            constraints.append(content)
        elif any(k in low for k in GOAL_HINTS):
            goals.append(content)
        elif b or any(k in low for k in FEATURE_HINTS):
            features.append(content)
    return {
        "title": title or "Untitled document",
        "goals": goals, "features": features,
        "constraints": constraints, "milestones": milestones,
    }


def _slug(text: str) -> str:
    s = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    return (s[:40] or "doc-intake").upper()


def propose_plan(analysis: dict, *, source: str = "") -> dict:
    """Build a B-mode taskset proposal (work.py-new shaped). Owner-gated; not registered."""
    ts_id = f"TASKSET-AR-{_slug(analysis['title'])}"
    features = analysis["features"] or analysis["goals"] or [analysis["title"]]
    tasks = [{
        "title": feat[:80],
        "summary": feat,
        "context": feat,
        "scope": "Derived from document intake; refine before dispatch.",
        "acceptance": ["Implements: " + feat[:120]],
    } for feat in features]
    return {
        "mode": "B-mode",
        "approval": "owner_gated",
        "origin_type": "doc_intake",
        "status": "proposed",
        "source_document": source,
        "taskset": {
            "task_set_id": ts_id,
            "display_name": analysis["title"][:60],
            "summary": " ".join(analysis["goals"])[:400] or analysis["title"],
        },
        "tasks": tasks,
        "constraints": analysis["constraints"],
        "milestones": analysis["milestones"],
        "note": "Submit via `work.py new` ONLY after Owner approval; this is a proposal, "
                "not a registration.",
    }


def doc_to_plan(path: Path, *, title: str | None = None) -> dict:
    analysis = analyze(extract_text(path), title=title)
    return propose_plan(analysis, source=str(path))


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Document -> registerable plan proposal (Owner-gated).")
    ap.add_argument("--input", required=True)
    ap.add_argument("--title")
    ap.add_argument("--out", help="write the proposal JSON here for Owner review")
    a = ap.parse_args(argv)
    proposal = doc_to_plan(Path(a.input), title=a.title)
    text = json.dumps(proposal, indent=2, ensure_ascii=False)
    if a.out:
        Path(a.out).write_text(text + "\n", encoding="utf-8")
        print(f"doc-to-plan: proposal written to {a.out} ({len(proposal['tasks'])} tasks) - Owner approval required")
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
